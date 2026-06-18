#!/usr/bin/env python3
"""
Title: Convening Slide Deck Builder - Native PowerPoint Objects

Description:
    Builds the same "Happening Now / Up Next" registration deck as
    build_deck.py, but using fully editable, native PowerPoint objects
    (rectangles, text boxes, and picture placeholders) instead of one
    flattened PNG per slide.

    Everything is editable in PowerPoint: titles and labels are real text,
    the stripe bar / tabs / bands are real shapes, and each headshot is a
    picture object you can swap with right-click -> Change Picture.

    Layout, palette, fonts, and agenda data are imported from build_deck.py
    so the two outputs stay in sync. Canvas coordinates (2560 x 1440 px) are
    mapped onto the 13.333 x 7.5 in (16:9) slide; font sizes convert at
    0.375 pt per px.

    NOTE ON FONTS: native text renders in whatever fonts the presenting
    machine has installed. Install Oswald + Public Sans (see deck/fonts/) or
    use PowerPoint's "Embed fonts in the file" save option to lock the look;
    otherwise PowerPoint substitutes a default sans-serif.

Usage:
    python build_deck_native.py        # writes convening-registration-deck-editable.pptx
    python build_deck_native.py --verify   # also print a structural summary

Changelog:
    2026-06-15  Initial native-object version.
"""

import argparse
import os

from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

import build_deck as bd

# --------------------------------------------------------------------------
# Coordinate / unit mapping (canvas px -> slide EMU and pt)
# --------------------------------------------------------------------------
WIDTH, HEIGHT = bd.WIDTH, bd.HEIGHT          # 2560 x 1440 design canvas
SLIDE_W_EMU = 12192000                        # 13.333 in
SLIDE_H_EMU = 6858000                         # 7.5 in
EMU_PER_PX = SLIDE_W_EMU / WIDTH              # uniform (matches vertical too)
PT_PER_PX = 0.375                             # 1440 px == 7.5 in == 540 pt

# Font family names as PowerPoint will look them up.
OSWALD = "Oswald"
PUBLIC_SANS = "Public Sans"

# Scratch drawing context used only to measure/wrap text with the real TTFs.
_SCRATCH = ImageDraw.Draw(Image.new("RGB", (10, 10)))

# Cached placeholder image path for speakers without a photo yet.
_PLACEHOLDER_PATH = os.path.join(bd.BUILD_DIR, "_placeholder.png")

# Cached blank placeholder (background only, no silhouette) for a speaker whose
# headshot is intentionally left empty.
_BLANK_PATH = os.path.join(bd.BUILD_DIR, "_blank.png")


def emu(px):
    """Convert a canvas pixel measurement to EMU for slide geometry."""
    return Emu(int(round(px * EMU_PER_PX)))


def pt(px):
    """Convert a canvas pixel font size to a PowerPoint point size."""
    return Pt(px * PT_PER_PX)


def color(rgb):
    """Wrap an (r, g, b) tuple as a python-pptx RGBColor."""
    return RGBColor(*rgb)


# --------------------------------------------------------------------------
# Low-level shape / text helpers
# --------------------------------------------------------------------------
def add_rect(slide, x, y, w, h, fill):
    """Add a borderless, shadowless filled rectangle in canvas coordinates."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, emu(x), emu(y), emu(w), emu(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color(fill)
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_text(slide, x, y, w, h, text=None, *, font_name, size_px, fill,
             bold=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             tracking_px=0, wrap=True, lines=None, line_spacing=None):
    """
    Add a text box with one paragraph per line.

    Parameters mirror the canvas renderer: positions/sizes are in canvas px,
    ``font_name`` is a PowerPoint family, ``size_px`` is a canvas font size,
    ``tracking_px`` adds letter spacing, and ``lines`` (if given) overrides
    ``text`` to place explicit, pre-wrapped lines.
    """
    box = slide.shapes.add_textbox(emu(x), emu(y), emu(w), emu(h))
    frame = box.text_frame
    frame.word_wrap = wrap
    frame.margin_left = frame.margin_right = 0
    frame.margin_top = frame.margin_bottom = 0
    frame.vertical_anchor = anchor

    paragraph_lines = lines if lines is not None else [text or ""]
    for index, line in enumerate(paragraph_lines):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.alignment = align
        if line_spacing:
            paragraph.line_spacing = line_spacing
        run = paragraph.add_run()
        run.text = line
        run.font.name = font_name
        run.font.size = pt(size_px)
        run.font.bold = bold
        run.font.color.rgb = color(fill)
        if tracking_px:
            # OOXML rPr@spc is in 1/100 pt; px -> pt is PT_PER_PX.
            run._r.get_or_add_rPr().set("spc", str(int(tracking_px * PT_PER_PX * 100)))
    return box


def add_tab(slide, x, y, text, *, font_key, font_name, size_px, fg, bg,
            tracking_px=8, pad_x=40, pad_y=22, anchor="left", bold=False):
    """
    Add a rectangular tab (filled block + centered label), sized to the text
    using the real TTF metrics so it matches the image renderer.
    """
    measure_font = bd.font(font_key, size_px)
    ascent, descent = measure_font.getmetrics()
    text_w = bd.measure_tracked(text, measure_font, tracking_px)
    box_w = text_w + pad_x * 2
    box_h = (ascent + descent) + pad_y * 2
    left = x if anchor == "left" else x - box_w
    add_rect(slide, left, y, box_w, box_h, bg)
    add_text(slide, left, y, box_w, box_h, text, font_name=font_name,
             size_px=size_px, fill=fg, bold=bold, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE, tracking_px=tracking_px, wrap=False)
    return left, box_w, box_h


def _set_oval_geometry(picture):
    """Mask a picture to an ellipse by replacing its preset geometry."""
    sppr = picture._element.spPr
    for tag in ("a:prstGeom", "a:custGeom"):
        existing = sppr.find(qn(tag))
        if existing is not None:
            sppr.remove(existing)
    geom = sppr.makeelement(qn("a:prstGeom"), {"prst": "ellipse"})
    geom.append(sppr.makeelement(qn("a:avLst"), {}))
    sppr.append(geom)


def add_headshot(slide, path, x, y, size, shape, blank=False):
    """
    Add a headshot as a native picture: square cover-cropped (top-biased) and,
    for ``shape == "circle"``, masked to a circle. Missing photos use the
    shared placeholder image so a real photo can be dropped in later; when
    ``blank`` is set, an empty placeholder disc (no silhouette) is used instead.
    """
    placeholder = _BLANK_PATH if blank else _PLACEHOLDER_PATH
    source = path if (path and os.path.exists(path)) else placeholder
    picture = slide.shapes.add_picture(source, emu(x), emu(y), emu(size), emu(size))

    img_w, img_h = Image.open(source).size
    if img_w > img_h:                                   # landscape: trim sides
        trim = (img_w - img_h) / img_w
        picture.crop_left = trim / 2
        picture.crop_right = trim / 2
    elif img_h > img_w:                                 # portrait: trim top/bottom
        trim = (img_h - img_w) / img_h
        picture.crop_top = trim * bd.HEAD_CROP_BIAS
        picture.crop_bottom = trim * (1 - bd.HEAD_CROP_BIAS)

    if shape != "square":
        _set_oval_geometry(picture)
    return picture


# --------------------------------------------------------------------------
# Slide regions
# --------------------------------------------------------------------------
def draw_background(slide):
    """Fill the slide cream and lay the black/gold/red/light stripe bar."""
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color(bd.PAPER)
    bar_h = 26
    segments = [
        (0.00, 0.22, bd.BLACK),
        (0.22, 0.46, bd.GOLD),
        (0.46, 0.74, bd.RED),
        (0.74, 1.00, bd.LIGHT),
    ]
    for start, end, seg_color in segments:
        add_rect(slide, WIDTH * start, 0, WIDTH * (end - start), bar_h, seg_color)


def place_speakers(slide, speakers, zone_top, zone_bottom, default_shape):
    """
    Lay out native headshots as a two-up grid with name + role + title beside
    each photo, matching the image renderer's geometry.

    Parameters:
        slide: the python-pptx slide to populate.
        speakers (list): (name, role, image filename[, mode]) tuples.
        zone_top/zone_bottom (int): vertical band the grid is centered within.
        default_shape (str): session default headshot shape ("circle"/"square").
    """
    count = len(speakers)
    if count == 0:
        return

    grid = bd.speaker_grid_cells(count, zone_top, zone_bottom)
    head = grid["head"]
    inner_gap = grid["inner_gap"]
    name_font = bd.font("sans-700", grid["name_size"])
    name_lh = int(grid["name_size"] * 1.12)
    role_lh = int(grid["role_size"] * 1.25)

    normalized = [
        (sp[0], sp[1], sp[2], sp[3] if len(sp) > 3 else default_shape)
        for sp in speakers
    ]

    for (name, role, filename, mode), (cell_x, cell_y, cell_w, cell_h) in zip(
        normalized, grid["cells"]
    ):
        path = os.path.join(bd.SPEAKER_IMG_DIR, filename) if filename else None
        text_x = cell_x + head + inner_gap
        text_w = cell_w - head - inner_gap

        name_lines = bd.wrap_text(_SCRATCH, name, name_font, text_w)[:2]
        title = bd.speaker_title_for(name)
        title_lines, title_size = [], None
        if title:
            _, title_lines, title_size = bd.fit_display_text(
                _SCRATCH, title, text_w, 2, grid["title_sizes"], font_name="sans-400"
            )
            title_lines = title_lines[:2]  # fit_display_text may exceed 2 at min size
        title_lh = int(title_size * 1.18) if title_size else 0

        block_h = len(name_lines) * name_lh + grid["gap_after_name"] + role_lh
        if title_lines:
            block_h += grid["gap_after_role"] + len(title_lines) * title_lh

        # Headshot, vertically centered in the cell. The image deck's "full"
        # (uncropped contain-fit) mode has no python-pptx equivalent, so any
        # non-square mode renders as a circle here.
        head_y = cell_y + (cell_h - head) // 2
        add_headshot(slide, path, cell_x, head_y, head,
                     "square" if mode == "square" else "circle",
                     blank=(mode == "blank"))

        # Text block (name, role, title), vertically centered against the photo.
        text_top = cell_y + (cell_h - block_h) // 2
        add_text(slide, text_x, text_top, text_w, name_lh * len(name_lines) + 8,
                 lines=name_lines, font_name=PUBLIC_SANS, size_px=grid["name_size"],
                 fill=bd.INK, bold=True, align=PP_ALIGN.LEFT)

        role_y = text_top + len(name_lines) * name_lh + grid["gap_after_name"]
        add_text(slide, text_x, role_y, text_w, role_lh + 6, role.upper(),
                 font_name=PUBLIC_SANS, size_px=grid["role_size"], fill=bd.RED_DARK,
                 bold=True, align=PP_ALIGN.LEFT, tracking_px=4)

        if title_lines:
            title_y = role_y + role_lh + grid["gap_after_role"]
            add_text(slide, text_x, title_y, text_w, title_lh * len(title_lines) + 8,
                     lines=title_lines, font_name=PUBLIC_SANS, size_px=title_size,
                     fill=bd.INK_SOFT, bold=False, align=PP_ALIGN.LEFT)


def place_note(slide, note, zone_top, zone_bottom):
    """Centered supporting line + small gold square (slots without speakers)."""
    if not note:
        return
    center_y = (zone_top + zone_bottom) // 2
    add_rect(slide, WIDTH / 2 - 17, center_y - 92, 34, 34, bd.GOLD)
    _, lines, size = bd.fit_display_text(
        _SCRATCH, note, WIDTH - bd.MARGIN * 2 - 200, 2, [76, 66, 58, 50],
        font_name="oswald-500",
    )
    line_h = int(size * 1.12)
    add_text(slide, bd.MARGIN, center_y - (len(lines) * line_h) // 2 + 16,
             WIDTH - bd.MARGIN * 2, line_h * len(lines) + 24, lines=lines,
             font_name=OSWALD, size_px=size, fill=bd.INK_SOFT, align=PP_ALIGN.CENTER)


def draw_up_next(slide, next_item):
    """Dark 'Up Next' band: gold label, next title, and start time."""
    band_top = 1208
    add_rect(slide, 0, band_top, WIDTH, HEIGHT - band_top, bd.DARK)
    add_rect(slide, 0, band_top, WIDTH, 7, bd.GOLD)
    add_text(slide, bd.MARGIN, band_top + 44, 700, 64, "UP NEXT", font_name=OSWALD,
             size_px=46, fill=bd.GOLD, bold=True, tracking_px=10)
    _, lines, size = bd.fit_display_text(
        _SCRATCH, next_item["title"], WIDTH - bd.MARGIN * 2 - 360, 2,
        [70, 62, 54, 48], font_name="sans-700",
    )
    add_text(slide, bd.MARGIN, band_top + 104, WIDTH - bd.MARGIN * 2 - 360,
             int(size * 1.1) * len(lines) + 20, lines=lines, font_name=PUBLIC_SANS,
             size_px=size, fill=bd.LIGHT, bold=True, line_spacing=1.05)
    add_text(slide, WIDTH - bd.MARGIN - 520, band_top + 84, 520, 110,
             next_item["start"], font_name=OSWALD, size_px=84, fill=bd.GOLD,
             bold=True, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)


# --------------------------------------------------------------------------
# Whole slides
# --------------------------------------------------------------------------
def build_slot_slide(slide, item, next_item):
    """Populate one 'Happening Now / Up Next' slide with native objects."""
    draw_background(slide)

    add_tab(slide, bd.MARGIN, 132, "HAPPENING NOW", font_key="oswald-600",
            font_name=OSWALD, size_px=44, fg=bd.LIGHT, bg=bd.RED, tracking_px=8, bold=True)
    add_tab(slide, WIDTH - bd.MARGIN, 132, item["time"], font_key="oswald-600",
            font_name=OSWALD, size_px=46, fg=bd.INK, bg=bd.GOLD, tracking_px=4,
            anchor="right", bold=True)

    kicker, main_title = bd.split_kicker(item["title"])
    if kicker:
        add_text(slide, bd.MARGIN, 268, WIDTH - bd.MARGIN * 2, 72, kicker,
                 font_name=OSWALD, size_px=52, fill=bd.GOLD_DEEP, bold=True, tracking_px=10)
    title_font, lines, size = bd.fit_display_text(
        _SCRATCH, main_title, WIDTH - bd.MARGIN * 2, 2,
        [168, 154, 140, 126, 112, 100],
    )
    add_text(slide, bd.MARGIN, 318, WIDTH - bd.MARGIN * 2, int(size * 1.02) * 2 + 40,
             lines=lines, font_name=OSWALD, size_px=size, fill=bd.INK, bold=True,
             line_spacing=1.0)

    zone_top = max(318 + int(size * 1.02) * len(lines) + 40, 700)
    if item["speakers"]:
        place_speakers(slide, item["speakers"], zone_top, 1150,
                       item.get("headshot_shape", "circle"))
    else:
        place_note(slide, item.get("note", ""), zone_top, 1150)

    add_text(slide, bd.MARGIN, 1158, WIDTH / 2, 48, bd.EVENT_FOOTER,
             font_name=PUBLIC_SANS, size_px=30, fill=bd.INK_SOFT, bold=True)
    add_text(slide, WIDTH / 2, 1158, WIDTH / 2 - bd.MARGIN, 48, bd.EVENT_DATE,
             font_name=PUBLIC_SANS, size_px=30, fill=bd.INK_SOFT, bold=True,
             align=PP_ALIGN.RIGHT)

    draw_up_next(slide, next_item)


def build_cover_slide(slide):
    """Populate the opening welcome/cover slide with native objects."""
    draw_background(slide)

    add_tab(slide, bd.MARGIN, 200, bd.EVENT_TAGLINE.upper(), font_key="oswald-600",
            font_name=OSWALD, size_px=46, fg=bd.LIGHT, bg=bd.RED, tracking_px=10, bold=True)

    _, lines, size = bd.fit_display_text(
        _SCRATCH, bd.EVENT_TITLE, WIDTH - bd.MARGIN * 2, 3, [188, 172, 156, 140, 126],
    )
    line_h = int(size * 1.0)
    add_text(slide, bd.MARGIN, 330, WIDTH - bd.MARGIN * 2, line_h * len(lines) + 30,
             lines=lines, font_name=OSWALD, size_px=size, fill=bd.INK, bold=True,
             line_spacing=1.0)

    details_y = 330 + line_h * len(lines) + 78
    add_text(slide, bd.MARGIN, details_y, WIDTH - bd.MARGIN * 2, 80,
             f"{bd.EVENT_DATE}  -  {bd.EVENT_TIME}", font_name=PUBLIC_SANS,
             size_px=50, fill=bd.INK, bold=True)
    add_text(slide, bd.MARGIN, details_y + 78, WIDTH - bd.MARGIN * 2, 70,
             bd.EVENT_VENUE, font_name=PUBLIC_SANS, size_px=44, fill=bd.INK_SOFT)

    band_top = 1148
    add_rect(slide, 0, band_top, WIDTH, HEIGHT - band_top, bd.DARK)
    add_rect(slide, 0, band_top, WIDTH, 7, bd.GOLD)
    add_text(slide, bd.MARGIN, band_top + 48, 1400, 60, "WELCOME - PLEASE CHECK IN",
             font_name=OSWALD, size_px=44, fill=bd.GOLD, bold=True, tracking_px=8)
    add_text(slide, bd.MARGIN, band_top + 116, 1600, 60,
             "Find your name, grab a badge, and enjoy breakfast.",
             font_name=PUBLIC_SANS, size_px=40, fill=bd.LIGHT)

    # Partner logos on white cards, right-aligned within the band.
    card_h = 132
    card_y = band_top + (HEIGHT - band_top - card_h) // 2
    right = WIDTH - bd.MARGIN
    for filename in reversed(bd.PARTNER_LOGOS):
        path = os.path.join(bd.LOGO_IMG_DIR, filename)
        if not os.path.exists(path):
            continue
        logo = Image.open(path)
        max_w, max_h = 220, card_h - 36
        ratio = min(max_w / logo.width, max_h / logo.height)
        draw_w, draw_h = logo.width * ratio, logo.height * ratio
        card_w = draw_w + 48
        left = right - card_w
        add_rect(slide, left, card_y, card_w, card_h, bd.SURFACE)
        slide.shapes.add_picture(
            path, emu(left + (card_w - draw_w) / 2), emu(card_y + (card_h - draw_h) / 2),
            emu(draw_w), emu(draw_h),
        )
        right = left - 28


def build():
    """Build and save the native (editable) PowerPoint deck."""
    os.makedirs(bd.BUILD_DIR, exist_ok=True)
    bd.placeholder_face(800).save(_PLACEHOLDER_PATH)
    bd.blank_face(800).save(_BLANK_PATH)

    presentation = Presentation()
    presentation.slide_width = Emu(SLIDE_W_EMU)
    presentation.slide_height = Emu(SLIDE_H_EMU)
    blank_layout = presentation.slide_layouts[6]

    build_cover_slide(presentation.slides.add_slide(blank_layout))
    for item, next_item in bd.slide_items():
        build_slot_slide(presentation.slides.add_slide(blank_layout), item, next_item)

    output_path = os.path.join(bd.DECK_DIR, "convening-registration-deck-editable.pptx")
    presentation.save(output_path)
    return output_path, presentation


def main():
    parser = argparse.ArgumentParser(description="Build the editable native deck.")
    parser.add_argument("--verify", action="store_true",
                        help="Print a structural summary of the built deck.")
    args = parser.parse_args()

    output_path, presentation = build()
    print("Wrote", output_path)

    if args.verify:
        for index, slide in enumerate(presentation.slides):
            shapes = list(slide.shapes)
            pictures = sum(1 for s in shapes if s.shape_type == 13)
            texts = [s.text_frame.text.replace("\n", " / ") for s in shapes
                     if s.has_text_frame and s.text_frame.text.strip()]
            headline = texts[2] if len(texts) > 2 else (texts[0] if texts else "")
            print(f"  slide {index}: {len(shapes)} shapes, {pictures} pictures | {headline[:60]}")


if __name__ == "__main__":
    main()
