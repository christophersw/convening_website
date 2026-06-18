#!/usr/bin/env python3
"""
Title: Convening "Happening Now / Up Next" Slide Deck Builder

Description:
    Renders a striking, on-brand slide for each agenda time slot of the
    "Resilience Planning for Maryland Defense Communities" convening and
    assembles them into a 16:9 PowerPoint (.pptx) for display at the
    registration table.

    Each slot slide shows the current session as "Happening Now" (large
    title, time, and speaker headshots) and the following slot as
    "Up Next". A cover slide opens the deck. Slides are first rendered to
    high-resolution PNGs with Pillow (so typography and layout look
    identical on any computer), then placed full-bleed into the .pptx.

    Visual identity mirrors the event website: cream "paper" background,
    the black/gold/red/light stripe bar, Oswald condensed display type
    over Public Sans, brand red (#b01923) and gold (#f0b323).

Usage:
    python build_deck.py            # render every slide + build the .pptx
    python build_deck.py --only 7   # render a single slide (by index) only
    python build_deck.py --no-pptx  # render PNGs only, skip .pptx assembly

Changelog:
    2026-06-15  Initial version (cover + 12 slot slides, .pptx assembly).
"""

import argparse
import math
import os
import re

from PIL import Image, ImageDraw, ImageFont, ImageFilter

# --------------------------------------------------------------------------
# Paths
# --------------------------------------------------------------------------
DECK_DIR = os.path.dirname(os.path.abspath(__file__))
PROJECT_DIR = os.path.dirname(DECK_DIR)
FONT_DIR = os.path.join(DECK_DIR, "fonts")
BUILD_DIR = os.path.join(DECK_DIR, "build")
SPEAKER_IMG_DIR = os.path.join(PROJECT_DIR, "assets", "images", "speakers")
LOGO_IMG_DIR = os.path.join(PROJECT_DIR, "assets", "images", "logos")
SPEAKER_PAGES_DIR = os.path.join(PROJECT_DIR, "speakers")

# --------------------------------------------------------------------------
# Canvas + brand palette (RGB), taken from the site's CSS custom properties
# --------------------------------------------------------------------------
WIDTH, HEIGHT = 2560, 1440          # 16:9 at high resolution
MARGIN = 170

INK = (20, 19, 19)                  # --ink           #141313
INK_SOFT = (78, 63, 55)             # --ink-soft      #4e3f37
PAPER = (247, 243, 234)             # --paper         #f7f3ea
PAPER_2 = (239, 229, 213)           # --paper-2       #efe5d5
SURFACE = (255, 253, 247)           # --surface       #fffdf7
LINE = (185, 170, 150)              # --line          #b9aa96
RED = (176, 25, 35)                 # --brand-red     #b01923
RED_DARK = (121, 15, 22)            # deep red        #790f16
GOLD = (240, 179, 35)               # --brand-gold    #f0b323
GOLD_DEEP = (140, 106, 19)          # gold edge       #8c6a13
BLACK = (18, 18, 18)                # --brand-black   #121212
DARK = (23, 23, 23)                 # dark surface    #171717
LIGHT = (247, 240, 226)             # --brand-light   #f7f0e2
CREAM_RING = (252, 247, 236)

# Vertical crop bias for headshots: 0.0 keeps the very top of the photo,
# 0.5 is a centered crop. A low value keeps heads in frame on portrait photos.
HEAD_CROP_BIAS = 0.18

# --------------------------------------------------------------------------
# Fonts
# --------------------------------------------------------------------------
_FONT_FILES = {
    "oswald-700": "Oswald-700.ttf",
    "oswald-600": "Oswald-600.ttf",
    "oswald-500": "Oswald-500.ttf",
    "sans-700": "PublicSans-700.ttf",
    "sans-600": "PublicSans-600.ttf",
    "sans-400": "PublicSans-400.ttf",
}
_font_cache = {}


def font(name, size):
    """
    Load (and cache) a TrueType font.

    Parameters:
        name (str): a key from _FONT_FILES, e.g. "oswald-700".
        size (int): font size in pixels.

    Returns:
        PIL.ImageFont.FreeTypeFont: the requested font instance.
    """
    key = (name, size)
    if key not in _font_cache:
        path = os.path.join(FONT_DIR, _FONT_FILES[name])
        _font_cache[key] = ImageFont.truetype(path, size)
    return _font_cache[key]


# --------------------------------------------------------------------------
# Speaker job titles (sourced from the website speaker pages)
# --------------------------------------------------------------------------
_speaker_titles_cache = None


def slugify(name):
    """
    Convert a speaker's display name to their page slug.

    Lowercases the name, replaces every run of non-alphanumeric characters
    with a single hyphen, and strips leading/trailing hyphens, matching the
    speakers/<slug>.md filenames.

    Parameters:
        name (str): a speaker's display name, e.g. "LaRon Russell".

    Returns:
        str: the slug, e.g. "laron-russell".
    """
    return re.sub(r"[^a-z0-9]+", "-", name.lower()).strip("-")


def _read_front_matter_value(path, key):
    """
    Return a single value from a Markdown file's YAML front matter.

    Reads the block between the first two `---` fences and returns the value
    for ``key`` with any surrounding single/double quotes stripped.

    Parameters:
        path (str): path to the .md file.
        key (str): front-matter key, e.g. "speaker_title".

    Returns:
        str | None: the value, or None if the key is absent.
    """
    with open(path, encoding="utf-8") as handle:
        lines = handle.read().splitlines()
    if not lines or lines[0].strip() != "---":
        return None
    prefix = key + ":"
    for line in lines[1:]:
        if line.strip() == "---":
            break
        if line.startswith(prefix):
            value = line[len(prefix):].strip()
            if len(value) >= 2 and value[0] in "\"'" and value[-1] == value[0]:
                value = value[1:-1]
            return value
    return None


def load_speaker_titles():
    """
    Read job titles from the website speaker pages.

    Scans speakers/*.md, pulls the ``speaker_title`` value from each file's
    leading front-matter block, and caches the result for the build.

    Returns:
        dict[str, str]: {page slug: job title}, e.g.
        {"kim-pezza": "Climate Resilience Director, Maryland Comptroller's Office"}.
    """
    global _speaker_titles_cache
    if _speaker_titles_cache is not None:
        return _speaker_titles_cache

    titles = {}
    if os.path.isdir(SPEAKER_PAGES_DIR):
        for filename in sorted(os.listdir(SPEAKER_PAGES_DIR)):
            if not filename.endswith(".md"):
                continue
            title = _read_front_matter_value(
                os.path.join(SPEAKER_PAGES_DIR, filename), "speaker_title"
            )
            if title:
                titles[filename[:-3]] = title
    _speaker_titles_cache = titles
    return titles


def speaker_title_for(name):
    """
    Look up a speaker's job title by display name.

    Parameters:
        name (str): speaker display name.

    Returns:
        str | None: the job title from their website page, or None if there
        is no matching speakers/<slug>.md or no title in it.
    """
    return load_speaker_titles().get(slugify(name))


# --------------------------------------------------------------------------
# Agenda data
#
# The full day is listed in order (including Registration and Adjourn) so that
# every slide can reference the slot that follows it as "Up Next". Only the
# items whose ``slide`` flag is True are rendered as slides.
#
# Each item:
#   key        unique short id
#   title      session title (a "Prefix: Rest" title is split into a gold
#              kicker + main title automatically)
#   start      human start time used in the "Up Next" line
#   time       full time range shown on the "Happening Now" slide
#   kind       "session" | "break" | "meal"
#   slide      whether to render a dedicated slide for this slot
#   note       optional supporting line for slots without speakers
#   speakers   list of (name, role, image filename) tuples
# --------------------------------------------------------------------------
AGENDA = [
    {
        "key": "registration",
        "title": "Registration",
        "start": "7:30 AM",
        "time": "7:30 - 8:30 AM",
        "kind": "session",
        "slide": False,
        "speakers": [],
    },
    {
        "key": "welcoming-remarks",
        "title": "Welcoming Remarks",
        "start": "8:30 AM",
        "time": "8:30 - 8:45 AM",
        "kind": "session",
        "slide": True,
        "speakers": [
            ("Lisa Swoboda", "Speaker", "lisa-swoboda.png"),
        ],
    },
    {
        "key": "project-findings",
        "title": "Resilient Maryland Defense Communities Project Findings",
        "start": "8:45 AM",
        "time": "8:45 - 10:00 AM",
        "kind": "session",
        "slide": True,
        "speakers": [
            ("Chris Webster", "Presenter", "chris-webster.jpg"),
        ],
    },
    {
        "key": "fireside-chat",
        "title": "Fireside Chat: Flooding and Community Partnerships",
        "start": "10:00 AM",
        "time": "10:00 - 10:45 AM",
        "kind": "session",
        "slide": True,
        "speakers": [
            ("Cyrena Chiles Eitler", "Moderator", "cyrena-chiles-eitler.jpg"),
            ("Eric Leshinsky", "Panelist", None),  # photo to be added manually
        ],
    },
    {
        "key": "break-am",
        "title": "Break",
        "start": "10:45 AM",
        "time": "10:45 - 11:00 AM",
        "kind": "break",
        "slide": True,
        "note": "Refreshments & Networking",
        "speakers": [],
    },
    {
        "key": "energy-panel",
        "title": "Panel: Innovative Energy Partnerships",
        "start": "11:00 AM",
        "time": "11:00 - 11:45 AM",
        "kind": "session",
        "slide": True,
        "speakers": [
            ("Harry Kleiser", "Moderator", None),  # photo to be added manually
            ("Shawn Butler", "Panelist", "shawn-butler.png"),
            ("Brian Lazarchick", "Panelist", "brian-lazarchick.jpg"),
        ],
    },
    {
        "key": "lunch",
        "title": "Catered Lunch & Networking",
        "start": "11:45 AM",
        "time": "11:45 AM - 12:45 PM",
        "kind": "meal",
        "slide": True,
        "note": "Please enjoy lunch and connect with fellow attendees",
        "speakers": [],
    },
    {
        "key": "federal-state-financing",
        "title": "Panel: Federal and State Financing Opportunities",
        "start": "12:45 PM",
        "time": "12:45 - 1:35 PM",
        "kind": "session",
        "slide": True,
        "speakers": [
            ("Kim Pezza", "Moderator", "kim-pezza.jpg"),
            ("Patricia Gray", "Panelist", "patricia-gray.jpg"),
            ("Leah Sheppard", "Panelist", "leah-sheppard.jpg"),
            ("LaRon Russell", "Panelist", "laron-russell.png"),
        ],
    },
    {
        "key": "creative-financing",
        "title": "Panel: Creative Financing Opportunities",
        "start": "1:35 PM",
        "time": "1:35 - 2:25 PM",
        "kind": "session",
        "slide": True,
        "speakers": [
            ("Mark Belton", "Moderator", "mark-belton.jpg"),
            ("Annabel Cryan", "Panelist", "annabel-cryan.jpg"),
            ("Mary Edwards", "Panelist", "mary-edwards.jpg"),
            ("Amanda Mock", "Panelist", None, "blank"),  # no headshot — blank placeholder disc
        ],
    },
    {
        "key": "break-pm",
        "title": "Break",
        "start": "2:25 PM",
        "time": "2:25 - 2:35 PM",
        "kind": "break",
        "slide": True,
        "note": "Refreshments & Networking",
        "speakers": [],
    },
    {
        "key": "regional-workshops",
        "title": "Regional Workshops / Breakouts",
        "start": "2:35 PM",
        "time": "2:35 - 3:20 PM",
        "kind": "session",
        "slide": True,
        "note": "Facilitated regional breakout discussions",
        "speakers": [],
    },
    {
        "key": "transition-break",
        "title": "Transition & Break",
        "start": "3:20 PM",
        "time": "3:20 - 3:30 PM",
        "kind": "break",
        "slide": True,
        "note": "Please return to the main room",
        "speakers": [],
    },
    {
        "key": "wrap-up",
        "title": "Wrap-Up and Next Steps",
        "start": "3:30 PM",
        "time": "3:30 - 4:00 PM",
        "kind": "session",
        "slide": True,
        "note": "Closing reflections and next steps",
        "speakers": [],
    },
    {
        "key": "adjourn",
        "title": "Adjourn",
        "start": "4:00 PM",
        "time": "4:00 PM",
        "kind": "session",
        "slide": False,
        "speakers": [],
    },
]

EVENT_TITLE = "Resilience Planning for Maryland Defense Communities"
EVENT_TAGLINE = "Statewide Convening"
EVENT_DATE = "Tuesday, June 23, 2026"
EVENT_TIME = "8:30 AM - 4:00 PM"
EVENT_VENUE = "University of Maryland, Baltimore - SMC Campus Center"
EVENT_FOOTER = "Resilience Planning for Maryland Defense Communities"
PARTNER_LOGOS = [
    "maryland-commerce.jpg",
    "chhs.png",
    "mdem-mor.png",
    "mdp.png",
]


# --------------------------------------------------------------------------
# Drawing helpers
# --------------------------------------------------------------------------
def measure_tracked(text, fnt, tracking):
    """Return the pixel width of ``text`` drawn with letter ``tracking``."""
    if not text:
        return 0
    width = sum(fnt.getlength(ch) for ch in text)
    return width + tracking * (len(text) - 1)


def draw_tracked(draw, xy, text, fnt, fill, tracking, anchor="la"):
    """
    Draw text with manual letter-spacing (Pillow has no native tracking).

    Parameters:
        draw (ImageDraw): target drawing context.
        xy (tuple): (x, y) anchor point; y is the text top.
        text (str): the string to draw.
        fnt: the font to use.
        fill: text color.
        tracking (float): extra pixels inserted between characters.
        anchor (str): horizontal anchor - "la" left, "ma" center, "ra" right.

    Returns:
        float: the total drawn width.
    """
    total = measure_tracked(text, fnt, tracking)
    x, y = xy
    if anchor == "ra":
        x -= total
    elif anchor == "ma":
        x -= total / 2
    for ch in text:
        draw.text((x, y), ch, font=fnt, fill=fill)
        x += fnt.getlength(ch) + tracking
    return total


def wrap_text(draw, text, fnt, max_width):
    """Greedy word-wrap ``text`` to fit ``max_width``; returns a list of lines."""
    lines, current = [], ""
    for word in text.split():
        candidate = (current + " " + word).strip()
        if draw.textlength(candidate, font=fnt) <= max_width or not current:
            current = candidate
        else:
            lines.append(current)
            current = word
    if current:
        lines.append(current)
    return lines


def fit_display_text(draw, text, max_width, max_lines, sizes, font_name="oswald-700"):
    """
    Pick the largest font size from ``sizes`` (descending) at which ``text``
    wraps into at most ``max_lines`` lines within ``max_width``.

    Returns:
        tuple: (font, lines, size).
    """
    for size in sizes:
        fnt = font(font_name, size)
        lines = wrap_text(draw, text, fnt, max_width)
        if len(lines) <= max_lines:
            return fnt, lines, size
    fnt = font(font_name, sizes[-1])
    return fnt, wrap_text(draw, text, fnt, max_width), sizes[-1]


def tab(draw, x, y, text, fnt, *, fg, bg, tracking=8, pad_x=40, pad_y=22,
        anchor="left", radius=0):
    """
    Draw a rectangular tab/block containing tracked text.

    Parameters:
        x, y: anchor coordinates (top-left, or top-right when anchor="right").
        text: tab label.
        fnt: label font.
        fg/bg: text and fill colors.
        tracking: letter spacing for the label.
        pad_x/pad_y: horizontal/vertical padding inside the tab.
        anchor: "left" (x is left edge) or "right" (x is right edge).
        radius: corner radius in pixels (0 = crisp square corners).

    Returns:
        tuple: (left, top, right, bottom) bounding box of the tab.
    """
    ascent, descent = fnt.getmetrics()
    text_h = ascent + descent
    text_w = measure_tracked(text, fnt, tracking)
    box_w = text_w + pad_x * 2
    box_h = text_h + pad_y * 2
    left = x if anchor == "left" else x - box_w
    top = y
    right = left + box_w
    bottom = top + box_h
    if radius > 0:
        draw.rounded_rectangle((left, top, right, bottom), radius=radius, fill=bg)
    else:
        draw.rectangle((left, top, right, bottom), fill=bg)
    draw_tracked(draw, (left + pad_x, top + pad_y - descent * 0.15), text, fnt, fg, tracking)
    return left, top, right, bottom


def placeholder_face(size):
    """
    Build a neutral "photo goes here" silhouette tile, used when a speaker has
    no photo yet (so a headshot can be dropped in later).

    Parameters:
        size (int): edge length in pixels.

    Returns:
        PIL.Image: an RGB image of size (size, size).
    """
    img = Image.new("RGB", (size, size), (231, 223, 210))
    draw = ImageDraw.Draw(img)
    silhouette = (199, 187, 170)
    head_r = size * 0.155
    head_cx, head_cy = size / 2, size * 0.40
    draw.ellipse(
        (head_cx - head_r, head_cy - head_r, head_cx + head_r, head_cy + head_r),
        fill=silhouette,
    )
    draw.ellipse(
        (size * 0.17, size * 0.62, size * 0.83, size * 1.18), fill=silhouette
    )
    return img


def blank_face(size):
    """
    Build a blank placeholder tile: the same neutral background as
    ``placeholder_face`` but with no silhouette. Used for a speaker whose
    headshot is intentionally left empty (a plain empty disc rather than a
    "photo goes here" silhouette icon).

    Parameters:
        size (int): edge length in pixels.

    Returns:
        PIL.Image: an RGB image of size (size, size).
    """
    return Image.new("RGB", (size, size), (231, 223, 210))


def headshot(path, size, shape="circle", blank=False):
    """
    Build a borderless RGBA headshot badge cropped to a circle or square.

    Rendered at 4x and downscaled so the circular edge stays smooth.

    Parameters:
        path (str): image file path; a missing file yields a neutral disc.
        size (int): output edge length / diameter in pixels.
        shape (str): "circle" or "square".
        blank (bool): when no photo is available, render an empty placeholder
            disc (no silhouette) instead of the "photo goes here" silhouette.

    Returns:
        PIL.Image: an RGBA image of size (size, size).
    """
    scale = 4
    big = size * scale
    if path and os.path.exists(path):
        img = Image.open(path).convert("RGB")
        side = min(img.size)
        left = (img.width - side) // 2
        # Bias the vertical crop upward so the top of the head stays in frame.
        top = int((img.height - side) * HEAD_CROP_BIAS)
        img = img.crop((left, top, left + side, top + side)).resize(
            (big, big), Image.LANCZOS
        )
    elif blank:
        img = blank_face(big)
    else:
        img = placeholder_face(big)

    badge = img.convert("RGBA")
    if shape == "circle":
        mask = Image.new("L", (big, big), 0)
        ImageDraw.Draw(mask).ellipse((0, 0, big - 1, big - 1), fill=255)
        badge.putalpha(mask)
    return badge.resize((size, size), Image.LANCZOS)


def headshot_full(path, box):
    """
    Build an uncropped RGBA headshot: the whole photo scaled to fit (contain)
    within a ``box`` x ``box`` area, preserving its native aspect ratio.

    Parameters:
        path (str): image file path; a missing file yields a neutral square.
        box (int): the bounding box edge length in pixels.

    Returns:
        PIL.Image: an RGBA image no larger than (box, box) in either dimension.
    """
    if not (path and os.path.exists(path)):
        return Image.new("RGBA", (box, box), PAPER_2 + (255,))
    img = Image.open(path).convert("RGB")
    ratio = min(box / img.width, box / img.height)
    fitted = (max(1, round(img.width * ratio)), max(1, round(img.height * ratio)))
    return img.resize(fitted, Image.LANCZOS).convert("RGBA")


def paste_with_shadow(canvas, badge, x, y, blur=20, alpha=46, offset=12, shape="circle"):
    """Composite ``badge`` onto ``canvas`` at (x, y) over a soft, subtle shadow."""
    shadow = Image.new("RGBA", canvas.size, (0, 0, 0, 0))
    sdraw = ImageDraw.Draw(shadow)
    w, h = badge.size
    box = (x, y + offset, x + w, y + h + offset)
    if shape == "circle":
        sdraw.ellipse(box, fill=(40, 28, 20, alpha))
    else:
        sdraw.rectangle(box, fill=(40, 28, 20, alpha))
    shadow = shadow.filter(ImageFilter.GaussianBlur(blur))
    canvas.alpha_composite(shadow)
    canvas.alpha_composite(badge, (x, y))


def vignette(canvas, strength=0.05):
    """Darken the corners slightly for depth (subtle radial vignette)."""
    radial = Image.radial_gradient("L").resize(canvas.size)  # 0 center -> 255 edge
    alpha = radial.point(lambda v: int(v * strength))
    overlay = Image.new("RGBA", canvas.size, (15, 12, 8, 0))
    overlay.putalpha(alpha)
    canvas.alpha_composite(overlay)


def stripe_bar(draw, height=26):
    """Draw the site's signature black/gold/red/light stripe bar at the top."""
    segments = [
        (0.00, 0.22, BLACK),
        (0.22, 0.46, GOLD),
        (0.46, 0.74, RED),
        (0.74, 1.00, LIGHT),
    ]
    for start, end, color in segments:
        draw.rectangle((WIDTH * start, 0, WIDTH * end, height), fill=color)


def new_canvas():
    """Create the base cream canvas (RGBA) with stripe bar and vignette."""
    canvas = Image.new("RGBA", (WIDTH, HEIGHT), PAPER + (255,))
    vignette(canvas)
    draw = ImageDraw.Draw(canvas)
    stripe_bar(draw)
    return canvas, draw


def split_kicker(title):
    """
    Split a "Prefix: Rest" title into an uppercase kicker and main title.

    Returns:
        tuple: (kicker or None, main_title).
    """
    if ": " in title:
        prefix, rest = title.split(": ", 1)
        if len(prefix) <= 18:
            return prefix.upper(), rest
    return None, title


def speaker_grid_cells(count, zone_top, zone_bottom):
    """
    Compute the two-up grid geometry for a slide's speaker zone.

    Lays ``count`` speakers into a grid (1 -> one centered cell, otherwise two
    columns filled left-to-right, top-to-bottom) and vertically centers the
    grid within [zone_top, zone_bottom]. A final short row and a lone single
    cell are centered horizontally. The same geometry feeds both the image
    renderer and the native-object renderer so the two decks line up.

    Parameters:
        count (int): number of speakers (>= 1).
        zone_top (int): top of the available vertical band, in canvas px.
        zone_bottom (int): bottom of the band, in canvas px.

    Returns:
        dict: {
            "columns" (int), "rows" (int),
            "head" (int): headshot diameter in px,
            "inner_gap" (int): px between headshot and text block,
            "name_size" (int), "role_size" (int): canvas font sizes,
            "title_sizes" (list[int]): descending sizes for title auto-fit,
            "gap_after_name" (int), "gap_after_role" (int): px,
            "cells" (list[tuple]): one (cell_x, cell_y, cell_w, cell_h) per
                speaker, in speaker order.
        }
    """
    if not 1 <= count <= 4:
        raise ValueError(
            f"speaker_grid_cells supports 1-4 speakers, got {count}"
        )

    usable_w = WIDTH - MARGIN * 2
    col_gap = 90
    row_gap = 40

    if count == 1:
        # A lone speaker gets a fixed column narrower than usable_w so the
        # title text column is not absurdly wide across the whole slide.
        columns, col_w, head, cell_h = 1, 1400, 240, 300
        name_size, role_size, title_sizes = 46, 30, [40, 35, 31, 27]
    else:
        columns = 2
        col_w = (usable_w - col_gap) // 2
        if math.ceil(count / columns) == 1:           # exactly two speakers
            head, cell_h, name_size, role_size = 220, 280, 44, 30
            title_sizes = [38, 33, 29, 25]
        else:                                          # three or four speakers
            head, cell_h, name_size, role_size = 156, 196, 40, 26
            title_sizes = [32, 29, 26, 23]

    rows = math.ceil(count / columns)
    total_h = rows * cell_h + (rows - 1) * row_gap
    grid_top = zone_top + max(0, (zone_bottom - zone_top - total_h) // 2)

    cells = []
    for index in range(count):
        row = index // columns
        col = index % columns
        cells_in_row = min(columns, count - row * columns)
        row_w = cells_in_row * col_w + (cells_in_row - 1) * col_gap
        row_x0 = (WIDTH - row_w) // 2
        cell_x = row_x0 + col * (col_w + col_gap)
        cell_y = grid_top + row * (cell_h + row_gap)
        cells.append((cell_x, cell_y, col_w, cell_h))

    return {
        "columns": columns,
        "rows": rows,
        "head": head,
        "inner_gap": 28,
        "name_size": name_size,
        "role_size": role_size,
        "title_sizes": title_sizes,
        "gap_after_name": 10,
        "gap_after_role": 8,
        "cells": cells,
    }


# --------------------------------------------------------------------------
# Slide composition
# --------------------------------------------------------------------------
def draw_speaker_row(canvas, draw, speakers, zone_top, zone_bottom, shape="circle"):
    """
    Lay out speakers as a two-up grid: each headshot at left with the speaker's
    name, panel role, and job title stacked to its right.

    Parameters:
        canvas/draw: target image and its draw context.
        speakers (list): (name, role, image filename[, mode]) tuples.
        zone_top/zone_bottom (int): vertical band the grid is centered within.
        shape (str): default headshot shape, "circle" or "square".
    """
    count = len(speakers)
    if count == 0:
        return

    grid = speaker_grid_cells(count, zone_top, zone_bottom)
    head = grid["head"]
    inner_gap = grid["inner_gap"]
    name_font = font("sans-700", grid["name_size"])
    role_font = font("sans-600", grid["role_size"])
    name_lh = int(grid["name_size"] * 1.12)
    role_lh = int(grid["role_size"] * 1.25)

    # A 4th tuple element overrides the session shape; "full" shows the photo
    # uncropped, "square" / "circle" pick the badge shape.
    normalized = [
        (sp[0], sp[1], sp[2], sp[3] if len(sp) > 3 else shape) for sp in speakers
    ]

    for (name, role, filename, mode), (cell_x, cell_y, cell_w, cell_h) in zip(
        normalized, grid["cells"]
    ):
        path = os.path.join(SPEAKER_IMG_DIR, filename) if filename else None
        text_x = cell_x + head + inner_gap
        text_w = cell_w - head - inner_gap

        # Fit the name (<= 2 lines) and the job title (<= 2 lines, auto-shrunk).
        name_lines = wrap_text(draw, name, name_font, text_w)[:2]
        title = speaker_title_for(name)
        title_lines, title_size = [], None
        if title:
            _, title_lines, title_size = fit_display_text(
                draw, title, text_w, 2, grid["title_sizes"], font_name="sans-400"
            )
            title_lines = title_lines[:2]  # fit_display_text may exceed 2 at min size
        title_lh = int(title_size * 1.18) if title_size else 0

        block_h = len(name_lines) * name_lh + grid["gap_after_name"] + role_lh
        if title_lines:
            block_h += grid["gap_after_role"] + len(title_lines) * title_lh

        # Headshot, vertically centered in the cell.
        if mode == "full":
            badge = headshot_full(path, head)
        elif mode == "blank":
            badge = headshot(path, head, shape="circle", blank=True)
        else:
            badge = headshot(path, head, shape=mode)
        head_y = cell_y + (cell_h - head) // 2
        badge_x = cell_x + (head - badge.width) // 2
        badge_y = head_y + (head - badge.height) // 2
        paste_with_shadow(canvas, badge, badge_x, badge_y,
                          shape="circle" if mode in ("circle", "blank") else "square")

        # Text block, vertically centered against the headshot.
        ty = cell_y + (cell_h - block_h) // 2
        for line in name_lines:
            draw.text((text_x, ty), line, font=name_font, fill=INK, anchor="la")
            ty += name_lh
        ty += grid["gap_after_name"]
        draw_tracked(draw, (text_x, ty), role.upper(), role_font, RED_DARK, 4, anchor="la")
        ty += role_lh
        if title_lines:
            ty += grid["gap_after_role"]
            title_font = font("sans-400", title_size)
            for line in title_lines:
                draw.text((text_x, ty), line, font=title_font, fill=INK_SOFT, anchor="la")
                ty += title_lh


def draw_support_note(draw, note, zone_top, zone_bottom):
    """Draw a centered supporting line (used for slots without speakers)."""
    if not note:
        return
    center_y = (zone_top + zone_bottom) // 2
    # A small gold square above the note for a touch of structure.
    draw.rectangle((WIDTH / 2 - 17, center_y - 92, WIDTH / 2 + 17, center_y - 58), fill=GOLD)
    note_font, lines, _ = fit_display_text(
        draw, note, WIDTH - MARGIN * 2 - 200, 2, [76, 66, 58, 50], font_name="oswald-500"
    )
    lh = int(note_font.size * 1.12)
    ty = center_y - (len(lines) * lh) // 2 + 20
    for line in lines:
        draw.text((WIDTH / 2, ty), line, font=note_font, fill=INK_SOFT, anchor="ma")
        ty += lh


def draw_up_next(canvas, draw, next_item):
    """Draw the full-bleed dark 'Up Next' band at the bottom of the slide."""
    band_top = 1208
    draw.rectangle((0, band_top, WIDTH, HEIGHT), fill=DARK)
    draw.rectangle((0, band_top, WIDTH, band_top + 7), fill=GOLD)  # gold edge

    # Left: "UP NEXT" label + next session title and start time.
    label_font = font("oswald-600", 46)
    draw_tracked(draw, (MARGIN, band_top + 48), "UP NEXT", label_font, GOLD, 10)

    next_font, lines, _ = fit_display_text(
        draw,
        next_item["title"],
        WIDTH - MARGIN * 2 - 60,
        2,
        [70, 62, 54, 48],
        font_name="sans-700",
    )
    ty = band_top + 104
    lh = int(next_font.size * 1.1)
    for line in lines[:2]:
        draw.text((MARGIN, ty), line, font=next_font, fill=LIGHT, anchor="la")
        ty += lh

    # Right: the next slot's start time, large in gold.
    time_font = font("oswald-600", 84)
    draw.text((WIDTH - MARGIN, band_top + 96), next_item["start"], font=time_font,
              fill=GOLD, anchor="ra")


def render_slot_slide(item, next_item):
    """
    Render one 'Happening Now / Up Next' slide.

    Parameters:
        item (dict): the agenda slot being highlighted.
        next_item (dict): the slot that follows (shown as Up Next).

    Returns:
        PIL.Image: the finished RGB slide at WIDTH x HEIGHT.
    """
    canvas, draw = new_canvas()

    # --- Header row: "Happening Now" tab (left) + time tab (right) -----------
    tab_font = font("oswald-600", 44)
    tab(draw, MARGIN, 132, "HAPPENING NOW", tab_font, fg=LIGHT, bg=RED, tracking=8)
    time_font = font("oswald-600", 46)
    tab(draw, WIDTH - MARGIN, 132, item["time"], time_font, fg=INK, bg=GOLD,
        tracking=4, anchor="right")

    # --- Title block (optional gold kicker + large title) --------------------
    kicker, main_title = split_kicker(item["title"])
    title_top = 318
    if kicker:
        draw_tracked(draw, (MARGIN, 268), kicker, font("oswald-600", 52), GOLD_DEEP, 10)
    title_font, lines, size = fit_display_text(
        draw, main_title, WIDTH - MARGIN * 2, 2, [168, 154, 140, 126, 112, 100]
    )
    ty = title_top
    lh = int(size * 1.02)
    for line in lines[:2]:
        draw.text((MARGIN, ty), line, font=title_font, fill=INK, anchor="la")
        ty += lh
    title_bottom = ty

    # --- Middle zone: speakers, or a supporting note -------------------------
    zone_top = max(title_bottom + 40, 700)
    zone_bottom = 1150
    if item["speakers"]:
        draw_speaker_row(canvas, draw, item["speakers"], zone_top, zone_bottom,
                         shape=item.get("headshot_shape", "circle"))
    else:
        draw_support_note(draw, item.get("note", ""), zone_top, zone_bottom)

    # --- Footer line above the Up Next band ----------------------------------
    foot_font = font("sans-600", 30)
    draw.text((MARGIN, 1158), EVENT_FOOTER, font=foot_font, fill=INK_SOFT, anchor="la")
    draw.text((WIDTH - MARGIN, 1158), EVENT_DATE, font=foot_font, fill=INK_SOFT, anchor="ra")

    # --- Up Next band --------------------------------------------------------
    draw_up_next(canvas, draw, next_item)

    return canvas.convert("RGB")


def render_cover_slide():
    """Render the opening cover slide (event title, details, partner logos)."""
    canvas, draw = new_canvas()

    # Eyebrow tab
    tab(draw, MARGIN, 200, EVENT_TAGLINE.upper(), font("oswald-600", 46),
        fg=LIGHT, bg=RED, tracking=10)

    # Big event title
    title_font, lines, size = fit_display_text(
        draw, EVENT_TITLE, WIDTH - MARGIN * 2, 3, [188, 172, 156, 140, 126]
    )
    ty = 330
    lh = int(size * 1.0)
    for line in lines:
        draw.text((MARGIN, ty), line, font=title_font, fill=INK, anchor="la")
        ty += lh

    # Details
    ty += 78
    detail_font = font("sans-600", 50)
    draw.text((MARGIN, ty), f"{EVENT_DATE}  -  {EVENT_TIME}", font=detail_font,
              fill=INK, anchor="la")
    ty += 78
    draw.text((MARGIN, ty), EVENT_VENUE, font=font("sans-400", 44), fill=INK_SOFT,
              anchor="la")

    # Welcome / check-in band at the bottom with partner logos
    band_top = 1148
    draw.rectangle((0, band_top, WIDTH, HEIGHT), fill=DARK)
    draw.rectangle((0, band_top, WIDTH, band_top + 7), fill=GOLD)
    draw_tracked(draw, (MARGIN, band_top + 52), "WELCOME - PLEASE CHECK IN",
                 font("oswald-600", 44), GOLD, 8)
    draw.text((MARGIN, band_top + 118),
              "Find your name, grab a badge, and enjoy breakfast.",
              font=font("sans-400", 40), fill=LIGHT, anchor="la")

    # Partner logos on white cards, right-aligned within the band
    card_h = 132
    card_y = band_top + (HEIGHT - band_top - card_h) // 2
    x = WIDTH - MARGIN
    for filename in reversed(PARTNER_LOGOS):
        path = os.path.join(LOGO_IMG_DIR, filename)
        if not os.path.exists(path):
            continue
        logo = Image.open(path).convert("RGBA")
        max_w, max_h = 220, card_h - 36
        ratio = min(max_w / logo.width, max_h / logo.height)
        logo = logo.resize((int(logo.width * ratio), int(logo.height * ratio)), Image.LANCZOS)
        card_w = logo.width + 48
        left = x - card_w
        draw.rectangle((left, card_y, x, card_y + card_h), fill=SURFACE)
        canvas.alpha_composite(
            logo, (left + (card_w - logo.width) // 2, card_y + (card_h - logo.height) // 2)
        )
        x = left - 28

    return canvas.convert("RGB")


# --------------------------------------------------------------------------
# Deck assembly
# --------------------------------------------------------------------------
def slide_items():
    """Return (item, next_item) pairs for every slot flagged for a slide."""
    pairs = []
    for index, item in enumerate(AGENDA):
        if not item["slide"]:
            continue
        next_item = AGENDA[index + 1] if index + 1 < len(AGENDA) else None
        pairs.append((item, next_item))
    return pairs


def build_pptx(png_paths, output_path):
    """Assemble full-bleed PNG slides into a 16:9 .pptx at ``output_path``."""
    from pptx import Presentation
    from pptx.util import Inches

    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)
    blank_layout = presentation.slide_layouts[6]
    for path in png_paths:
        slide = presentation.slides.add_slide(blank_layout)
        slide.shapes.add_picture(path, 0, 0, width=presentation.slide_width,
                                 height=presentation.slide_height)
    presentation.save(output_path)


def main():
    parser = argparse.ArgumentParser(description="Build the convening slide deck.")
    parser.add_argument("--only", type=int, default=None,
                        help="Render a single slot slide by AGENDA index, no .pptx.")
    parser.add_argument("--no-pptx", action="store_true",
                        help="Render PNGs only; skip .pptx assembly.")
    args = parser.parse_args()

    os.makedirs(BUILD_DIR, exist_ok=True)

    if args.only is not None:
        item = AGENDA[args.only]
        next_item = AGENDA[args.only + 1] if args.only + 1 < len(AGENDA) else None
        image = render_slot_slide(item, next_item)
        out = os.path.join(BUILD_DIR, f"sample-{item['key']}.png")
        image.save(out)
        print("Wrote", out)
        return

    png_paths = []

    cover_path = os.path.join(BUILD_DIR, "slide-00-cover.png")
    render_cover_slide().save(cover_path)
    png_paths.append(cover_path)
    print("Wrote", cover_path)

    for number, (item, next_item) in enumerate(slide_items(), start=1):
        out = os.path.join(BUILD_DIR, f"slide-{number:02d}-{item['key']}.png")
        render_slot_slide(item, next_item).save(out)
        png_paths.append(out)
        print("Wrote", out)

    if not args.no_pptx:
        pptx_path = os.path.join(DECK_DIR, "convening-registration-deck.pptx")
        build_pptx(png_paths, pptx_path)
        print("Wrote", pptx_path)


if __name__ == "__main__":
    main()
